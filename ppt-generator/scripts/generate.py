#!/usr/bin/env python3
"""
PPT 生成脚本
依赖: python-pptx
"""

import argparse
import os
import sys
from pathlib import Path
from datetime import datetime
import subprocess


def ensure_pptx():
    """确保 python-pptx 已安装"""
    try:
        from pptx import Presentation
    except ImportError:
        print("正在安装 python-pptx...")
        subprocess.run([sys.executable, "-m", "pip", "install", "python-pptx", "-q"], check=True)


def generate_ppt(topic, pages=10, style="business"):
    """生成 PPT"""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RgbColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    
    prs = Presentation()
    
    # 设置幻灯片尺寸为 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # 样式配置
    styles = {
        "business": {"bg": RgbColor(0x1a, 0x1a, 0x2e), "accent": RgbColor(0x00, 0xd4, 0xff)},
        "academic": {"bg": RgbColor(0xff, 0xff, 0xff), "accent": RgbColor(0x33, 0x66, 0x99)},
        "simple": {"bg": RgbColor(0xf5, 0xf5, 0xf5), "accent": RgbColor(0x33, 0x33, 0x33)}
    }
    
    s = styles.get(style, styles["business"])
    
    # 封面
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)
    
    # 添加背景色
    background = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Inches(0), Inches(0),
        prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = s["bg"]
    background.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = topic
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = s["accent"]
    p.alignment = PP_ALIGN.CENTER
    
    # 副标题
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(12.333), Inches(0.8))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = datetime.now().strftime("%Y-%m-%d")
    p.font.size = Pt(20)
    p.font.color.rgb = RgbColor(0xaa, 0xaa, 0xaa)
    p.alignment = PP_ALIGN.CENTER
    
    # 生成内容页
    content_slides = [
        "概述",
        "背景介绍",
        "核心概念",
        "主要特点",
        "应用场景",
        "优势分析",
        "挑战与解决方案",
        "未来趋势",
        "总结与展望"
    ]
    
    for i, content in enumerate(content_slides[:pages-2], 1):
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        
        # 背景
        background = slide.shapes.add_shape(
            1, Inches(0), Inches(0),
            prs.slide_width, prs.slide_height
        )
        background.fill.solid()
        background.fill.fore_color.rgb = s["bg"]
        background.line.fill.background()
        
        # 页码
        page_num = slide.shapes.add_textbox(Inches(12.5), Inches(6.8), Inches(0.5), Inches(0.3))
        tf = page_num.text_frame
        p = tf.paragraphs[0]
        p.text = str(i + 1)
        p.font.size = Pt(14)
        p.font.color.rgb = s["accent"]
        
        # 标题
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"{i}. {content}"
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = s["accent"]
        
        # 内容区域
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(12.333), Inches(4.5))
        tf = content_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"[在此添加{content}的详细内容]"
        p.font.size = Pt(18)
        if style != "business":
            p.font.color.rgb = RgbColor(0x33, 0x33, 0x33)
        else:
            p.font.color.rgb = RgbColor(0xcc, 0xcc, 0xcc)
        
        # 添加要点
        for j in range(3):
            p = tf.add_paragraph()
            p.text = f"• 要点 {j+1}"
            p.font.size = Pt(16)
            p.space_before = Pt(12)
    
    # 结束页
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    background = slide.shapes.add_shape(
        1, Inches(0), Inches(0),
        prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = s["bg"]
    background.line.fill.background()
    
    end_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(12.333), Inches(1))
    tf = end_box.text_frame
    p = tf.paragraphs[0]
    p.text = "感谢观看"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = s["accent"]
    p.alignment = PP_ALIGN.CENTER
    
    # 保存
    filename = f"{datetime.now().strftime('%Y-%m-%d')}-{topic.replace(' ', '_')}.pptx"
    prs.save(filename)
    return filename


def main():
    parser = argparse.ArgumentParser(description="生成 PPT")
    parser.add_argument("topic", help="PPT 主题")
    parser.add_argument("--pages", type=int, default=10, help="页数")
    parser.add_argument("--style", default="business", choices=["business", "academic", "simple"], help="风格")
    
    args = parser.parse_args()
    
    ensure_pptx()
    
    print(f"正在生成关于 '{args.topic}' 的 PPT...")
    filename = generate_ppt(args.topic, args.pages, args.style)
    print(f"✅ PPT 已生成: {filename}")


if __name__ == "__main__":
    main()
