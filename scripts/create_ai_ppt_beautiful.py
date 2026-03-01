from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import os

# 创建演示文稿 - 宽屏16:9
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def set_slide_background(slide, color):
    """设置幻灯片背景色"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_slide(prs, title, subtitle):
    """添加标题页 - 带渐变效果"""
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, RGBColor(25, 25, 112))  # 深蓝背景
    
    # 添加装饰矩形
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(25, 25, 112)
    shape.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # 副标题
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(135, 206, 250)
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, content_list, accent_color=RGBColor(30, 144, 255)):
    """添加内容页 - 带视觉元素"""
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, RGBColor(245, 245, 250))  # 浅灰背景
    
    # 顶部装饰条
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.15))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent_color
    bar.line.fill.background()
    
    # 标题背景
    title_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.3), Inches(0.4), Inches(12.7), Inches(1))
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
    title_bg.line.color.rgb = RGBColor(220, 220, 220)
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.6), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(25, 25, 112)
    
    # 内容区域背景
    content_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(1.6), Inches(12.7), Inches(5.5))
    content_bg.fill.solid()
    content_bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
    content_bg.line.color.rgb = RGBColor(230, 230, 230)
    
    # 内容
    content_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.8), Inches(12.1), Inches(5.2))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, text in enumerate(content_list):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(50, 50, 50)
        p.space_before = Pt(12)
        p.space_after = Pt(6)
    
    return slide

def add_icon_slide(prs, title, items, accent_color=RGBColor(30, 144, 255)):
    """添加带图标的内容页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, RGBColor(245, 245, 250))
    
    # 顶部装饰条
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.15))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent_color
    bar.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(25, 25, 112)
    
    # 两列布局
    col_width = 6
    left_x = 0.4
    right_x = 6.9
    start_y = 1.4
    
    for i, (col_x, item) in enumerate([(left_x, items[0]), (right_x, items[1])] if len(items) == 2 else [(left_x, items[0])]):
        # 卡片背景
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(col_x), Inches(start_y), Inches(col_width), Inches(5.5))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(255, 255, 255)
        card.line.color.rgb = RGBColor(220, 220, 220)
        
        # 图标区域（用圆形代替）
        icon_bg = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(col_x + 2.3), Inches(start_y + 0.3), Inches(1.4), Inches(1.4))
        icon_bg.fill.solid()
        icon_bg.fill.fore_color.rgb = accent_color
        icon_bg.line.fill.background()
        
        # 内容
        content_box = slide.shapes.add_textbox(Inches(col_x + 0.3), Inches(start_y + 1.9), Inches(col_width - 0.6), Inches(3.4))
        tf = content_box.text_frame
        tf.word_wrap = True
        
        for j, text in enumerate(item):
            if j == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = text
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(50, 50, 50)
            p.space_before = Pt(8)
            p.space_after = Pt(4)
    
    return slide

# ==================== 开始创建幻灯片 ====================

# 1. 标题页
add_title_slide(prs, "人工智能 (AI) 简介", "从概念到应用的未来之旅")

# 2. 什么是AI
add_content_slide(prs, "什么是人工智能？", [
    "人工智能是让机器模拟人类智能的技术",
    "包括学习、推理、感知、理解语言等能力",
    "目标是让机器像人一样思考和行动",
    "1956年由约翰·麦卡锡首次提出概念",
    "如今已渗透到我们生活的方方面面"
])

# 3. AI的主要分支 - 两列布局
add_icon_slide(prs, "AI的主要分支", [
    [
        "机器学习 (Machine Learning)",
        "• 从数据中学习模式",
        "• 自动改进算法性能",
        "",
        "深度学习 (Deep Learning)",
        "• 基于神经网络",
        "• 处理复杂数据",
    ],
    [
        "自然语言处理 (NLP)",
        "• 理解和生成语言",
        "• ChatGPT就是NLP应用",
        "",
        "计算机视觉",
        "• 图像识别与分析",
        "• 人脸识别、自动驾驶",
    ]
])

# 4. AI的应用领域
add_content_slide(prs, "AI的应用领域", [
    "智能助手 - Siri, Alexa, Jarvis",
    "自动驾驶 - 特斯拉 Autopilot, Waymo",
    "医疗诊断 - 医学影像识别, 新药研发",
    "金融分析 - 智能投顾, 风险评估",
    "内容创作 - ChatGPT写作, Midjourney绘画",
    "工业制造 - 预测性维护, 质量控制"
], RGBColor(0, 150, 136))

# 5. AI的优势
add_content_slide(prs, "AI的核心优势", [
    "效率提升 - 7×24小时不间断工作，永不疲倦",
    "数据分析 - 处理海量数据，发现隐藏规律",
    "精准决策 - 基于数据做出客观判断",
    "个性化服务 - 根据用户习惯定制体验",
    "危险替代 - 执行高风险任务，保护人类",
    "持续进化 - 不断学习优化，越用越智能"
], RGBColor(76, 175, 80))

# 6. AI的挑战
add_content_slide(prs, "AI面临的挑战", [
    "数据隐私 - 如何保护用户数据安全",
    "算法偏见 - AI可能继承人类的偏见",
    "就业影响 - 部分工作可能被AI替代",
    "伦理边界 - AI决策的道德困境",
    "安全风险 - 深度伪造、自动化攻击",
    "监管难题 - 如何制定合适的法律法规"
], RGBColor(244, 67, 54))

# 7. 未来趋势 - 两列
add_icon_slide(prs, "AI的未来趋势", [
    [
        "通用人工智能 (AGI)",
        "• 接近人类水平的智能",
        "• 能处理各种任务",
        "",
        "多模态AI",
        "• 同时处理文本、图像、语音",
        "• 更自然的交互方式",
    ],
    [
        "AI Agent智能体",
        "• 能自主完成任务",
        "• 像Jarvis一样贴心",
        "",
        "边缘AI",
        "• 在设备端本地运行",
        "• 更好的隐私保护",
    ]
], RGBColor(156, 39, 176))

# 8. 如何开始
add_content_slide(prs, "如何开始接触AI？", [
    "1. 学习Python编程语言 - AI开发的基础",
    "2. 了解机器学习和深度学习基础概念",
    "3. 使用AI工具 - ChatGPT, Claude, Gemini",
    "4. 尝试AI应用开发 - 从简单项目开始",
    "5. 关注AI伦理和安全 - 负责任地使用AI",
    "6. 持续学习 - 跟上快速发展的技术"
], RGBColor(255, 152, 0))

# 9. 总结
add_content_slide(prs, "总结", [
    "AI正在深刻改变我们的工作和生活方式",
    "它带来巨大机遇，也带来挑战和考验",
    "关键是负责任地发展和使用AI技术",
    "未来属于会善用AI的人",
    "",
    "就像钢铁侠的Jarvis一样，",
    "AI可以成为我们最强大的伙伴！"
], RGBColor(63, 81, 181))

# 10. 结束页
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)
set_slide_background(slide, RGBColor(25, 25, 112))

# 结束标题
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(12.333), Inches(1.5))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "谢谢观看！"
p.font.size = Pt(60)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)
p.alignment = PP_ALIGN.CENTER

# 副标题
sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(12.333), Inches(1))
tf = sub_box.text_frame
p = tf.paragraphs[0]
p.text = "Questions & Discussion"
p.font.size = Pt(32)
p.font.color.rgb = RGBColor(135, 206, 250)
p.alignment = PP_ALIGN.CENTER

# 底部装饰线
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4), Inches(6), Inches(5.333), Inches(0.03))
line.fill.solid()
line.fill.fore_color.rgb = RGBColor(135, 206, 250)
line.line.fill.background()

# 保存文件
output_path = r"C:\Users\USER\.openclaw\media\outbound\AI_Introduction_Beautiful.pptx"
prs.save(output_path)
print(f"PPT已生成: {output_path}")
print("文件大小: {:.2f} KB".format(os.path.getsize(output_path) / 1024))
