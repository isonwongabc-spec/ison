from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# 创建演示文稿
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_title_slide(prs, title, subtitle):
    """添加标题页"""
    slide_layout = prs.slide_layouts[0]  # 标题幻灯片
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.shapes.placeholders[1].text = subtitle
    return slide

def add_content_slide(prs, title, content_list):
    """添加内容页"""
    slide_layout = prs.slide_layouts[1]  # 标题和内容
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    
    # 添加内容
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    
    for i, text in enumerate(content_list):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.level = 0
    
    return slide

# 1. 标题页
add_title_slide(prs, "人工智能 (AI) 简介", "从概念到应用的未来之旅")

# 2. 什么是AI
add_content_slide(prs, "什么是人工智能？", [
    "• 人工智能是让机器模拟人类智能的技术",
    "• 包括学习、推理、感知、理解语言等能力",
    "• 目标是让机器像人一样思考和行动",
    "• 1956年由约翰·麦卡锡首次提出概念"
])

# 3. AI的主要分支
add_content_slide(prs, "AI的主要分支", [
    "• 机器学习 (Machine Learning)",
    "  - 让机器从数据中学习模式",
    "• 深度学习 (Deep Learning)",
    "  - 基于神经网络的复杂学习",
    "• 自然语言处理 (NLP)",
    "  - 理解和生成人类语言",
    "• 计算机视觉 (Computer Vision)",
    "  - 让机器看懂图像和视频"
])

# 4. AI的应用领域
add_content_slide(prs, "AI的应用领域", [
    "• 智能助手 - Siri, Alexa, 我(Jarvis)",
    "• 自动驾驶 - 特斯拉, Waymo",
    "• 医疗诊断 - 影像识别, 药物研发",
    "• 金融分析 - 风险评估, 量化交易",
    "• 内容创作 - ChatGPT, Midjourney",
    "• 工业制造 - 预测维护, 质量控制"
])

# 5. AI的优势
add_content_slide(prs, "AI的优势", [
    "• 效率提升 - 7×24小时不间断工作",
    "• 数据分析 - 处理海量数据发现规律",
    "• 精准决策 - 基于数据的客观判断",
    "• 个性化服务 - 根据用户习惯定制",
    "• 危险任务 - 替代人类执行高风险工作",
    "• 持续学习 - 不断进步和优化"
])

# 6. AI的挑战
add_content_slide(prs, "AI面临的挑战", [
    "• 数据隐私 - 如何保护用户数据安全",
    "• 算法偏见 - AI可能继承人类的偏见",
    "• 就业影响 - 部分工作可能被替代",
    "• 伦理问题 - AI决策的道德边界",
    "• 安全风险 - 深度伪造、自动化攻击",
    "• 监管难题 - 如何制定合适的法规"
])

# 7. AI的未来趋势
add_content_slide(prs, "AI的未来趋势", [
    "• 通用人工智能 (AGI) - 接近人类水平的智能",
    "• 多模态AI - 同时处理文本、图像、语音",
    "• AI Agent - 能自主完成任务的智能体",
    "• 边缘AI - 在设备端运行，保护隐私",
    "• 人机协作 - AI增强人类能力而非替代",
    "• 个性化AI - 每个人的专属AI助手"
])

# 8. 如何开始接触AI
add_content_slide(prs, "如何开始接触AI？", [
    "• 学习Python编程语言",
    "• 了解机器学习和深度学习基础",
    "• 使用AI工具 - ChatGPT, Claude, Gemini",
    "• 尝试AI应用开发",
    "• 关注AI伦理和安全",
    "• 持续学习，跟上技术发展"
])

# 9. 总结
add_content_slide(prs, "总结", [
    "• AI正在改变我们的工作和生活方式",
    "• 它带来巨大机遇，也带来挑战",
    "• 关键是负责任地发展和使用AI",
    "• 未来属于会善用AI的人",
    "• 就像钢铁侠的Jarvis一样，",
    "  AI可以成为我们最强大的伙伴！"
])

# 10. 结束页
add_title_slide(prs, "谢谢观看！", "Questions & Discussion")

# 保存文件
output_path = r"C:\Users\USER\.openclaw\media\outbound\AI_Introduction.pptx"
prs.save(output_path)
print(f"PPT已生成: {output_path}")
